from pathlib import Path
from typing import Optional
import fitz
from docx import Document
from zipfile import BadZipFile
import csv
import openpyxl
import xlrd
from pptx import Presentation

def extract_txt(path: Path) -> Optional[str]:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return None

def extract_docx(path: Path) -> Optional[str]:
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except BadZipFile:
        return None
    except Exception:
        return None

def extract_pdf(path: Path) -> Optional[str]:
    try:
        text = []
        with fitz.open(path) as pdf:
            for page in pdf:
                text.append(page.get_text("text"))
        return "\n".join(text)
    except Exception:
        return None

def extract_csv(path: Path) -> Optional[str]:
    try:
        rows = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" ".join(row))
        return "\n".join(rows)
    except Exception:
        return None

def extract_xlsx(path: Path) -> Optional[str]:
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None)
                text.append(row_text)
        return "\n".join(text)
    except Exception:
        return None

def extract_xls(path: Path) -> Optional[str]:
    try:
        wb = xlrd.open_workbook(path)
        text = []
        for sheet in wb.sheets():
            for r in range(sheet.nrows):
                row = sheet.row_values(r)
                row_text = " ".join(str(c) for c in row if c not in (None, ""))
                text.append(row_text)
        return "\n".join(text)
    except Exception:
        return None

def extract_pptx(path: Path) -> Optional[str]:
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n".join(text)
    except Exception:
        return None

EXTRACTORS = {
    ".txt": extract_txt,
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".csv": extract_csv,
    ".xlsx": extract_xlsx,
    ".xls": extract_xls,
    ".pptx": extract_pptx,
}
