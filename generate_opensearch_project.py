from pathlib import Path

PROJECT_NAME = "opensearch_fastapi"
ROOT = Path.cwd() #/ PROJECT_NAME
FILES = {}

# =========================
# requirements.txt
# =========================
FILES["requirements.txt"] = r'''
fastapi
uvicorn
opensearch-py
python-docx
pymupdf
openpyxl
xlrd
python-pptx
'''

# =========================
# main.py
# =========================
FILES["main.py"] = r'''
from fastapi import FastAPI
from routes.indexing_routes import router as indexing_router
from routes.search_routes import router as search_router

app = FastAPI(title="OpenSearch File Search API")

app.include_router(indexing_router, prefix="/api")
app.include_router(search_router, prefix="/api")

@app.get("/")
def root():
    return {
        "status": "success",
        "message": "API running",
        "endpoints": ["/api/index-folder", "/api/search"]
    }
'''

# =========================
# config/settings.py
# =========================
FILES["config/settings.py"] = r'''
OPENSEARCH_HOST = "localhost"
OPENSEARCH_PORT = 9200
OPENSEARCH_INDEX = "documents"
'''

# =========================
# utils/response.py
# =========================
FILES["utils/response.py"] = r'''
def success_response(message, data=None):
    return {
        "status": "success",
        "message": message,
        "data": data or {}
    }

def failure_response(message, errors=None):
    return {
        "status": "failure",
        "message": message,
        "errors": errors
    }
'''

# =========================
# extractors/file_extractors.py
# =========================
FILES["extractors/file_extractors.py"] = r'''
from pathlib import Path
from typing import Optional
import fitz
from docx import Document
from zipfile import BadZipFile
import csv
import openpyxl
import xlrd
from pptx import Presentation

def extract_txt(path: Path) -> Optional[str]:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return None

def extract_docx(path: Path) -> Optional[str]:
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except BadZipFile:
        return None
    except Exception:
        return None

def extract_pdf(path: Path) -> Optional[str]:
    try:
        text = []
        with fitz.open(path) as pdf:
            for page in pdf:
                text.append(page.get_text("text"))
        return "\n".join(text)
    except Exception:
        return None

def extract_csv(path: Path) -> Optional[str]:
    try:
        rows = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" ".join(row))
        return "\n".join(rows)
    except Exception:
        return None

def extract_xlsx(path: Path) -> Optional[str]:
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None)
                text.append(row_text)
        return "\n".join(text)
    except Exception:
        return None

def extract_xls(path: Path) -> Optional[str]:
    try:
        wb = xlrd.open_workbook(path)
        text = []
        for sheet in wb.sheets():
            for r in range(sheet.nrows):
                row = sheet.row_values(r)
                row_text = " ".join(str(c) for c in row if c not in (None, ""))
                text.append(row_text)
        return "\n".join(text)
    except Exception:
        return None

def extract_pptx(path: Path) -> Optional[str]:
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n".join(text)
    except Exception:
        return None

EXTRACTORS = {
    ".txt": extract_txt,
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".csv": extract_csv,
    ".xlsx": extract_xlsx,
    ".xls": extract_xls,
    ".pptx": extract_pptx,
}
'''

# =========================
# opensearch_client/client.py
# =========================
FILES["opensearch_client/client.py"] = r'''
from opensearchpy import OpenSearch
from config.settings import OPENSEARCH_HOST, OPENSEARCH_PORT

def get_client():
    return OpenSearch(
        hosts=[{"host": OPENSEARCH_HOST, "port": OPENSEARCH_PORT}],
        use_ssl=False,
        verify_certs=False
    )
'''

# =========================
# opensearch_client/indexer.py
# =========================
FILES["opensearch_client/indexer.py"] = r'''
from pathlib import Path
from datetime import datetime
from extractors.file_extractors import EXTRACTORS

class OpenSearchIndexer:
    def __init__(self, client, index_name):
        self.client = client
        self.index_name = index_name
        self._ensure_index()

    def _ensure_index(self):
        if self.client.indices.exists(index=self.index_name):
            return
        self.client.indices.create(
            index=self.index_name,
            body={
                "mappings": {
                    "properties": {
                        "path": {"type": "keyword"},
                        "filename": {"type": "text"},
                        "filetype": {"type": "keyword"},
                        "modified": {"type": "date"},
                        "size_bytes": {"type": "long"},
                        "content": {"type": "text"}
                    }
                }
            }
        )

    def index_folder(self, folder: str):
        folder = Path(folder)
        count = 0
        for file in folder.rglob("*"):
            if not file.is_file():
                continue
            extractor = EXTRACTORS.get(file.suffix.lower())
            if not extractor:
                continue
            content = extractor(file)
            if not content:
                continue

            self.client.index(
                index=self.index_name,
                id=str(file.resolve()),
                body={
                    "path": str(file.resolve()),
                    "filename": file.name,
                    "filetype": file.suffix.lower().lstrip("."),
                    "modified": datetime.utcfromtimestamp(
                        file.stat().st_mtime
                    ).isoformat(),
                    "size_bytes": file.stat().st_size,
                    "content": content
                }
            )
            count += 1
        return count
'''

# =========================
# models/indexing_models.py
# =========================
FILES["models/indexing_models.py"] = r'''
from pydantic import BaseModel

class FolderInput(BaseModel):
    folder: str
'''

# =========================
# models/search_models.py
# =========================
FILES["models/search_models.py"] = r'''
from pydantic import BaseModel
from typing import List, Optional

class SearchInput(BaseModel):
    keyword: str
    search_mode: str  # filename or content
    file_types: Optional[List[str]] = None
    date_from: Optional[str] = None
    date_to: Optional[str] = None
    size_from: Optional[int] = None  # KB
    size_to: Optional[int] = None    # KB
    from_: Optional[int] = 0
    size: Optional[int] = 20
'''

# =========================
# routes/indexing_routes.py
# =========================
FILES["routes/indexing_routes.py"] = r'''
from fastapi import APIRouter, HTTPException
from models.indexing_models import FolderInput
from opensearch_client.client import get_client
from opensearch_client.indexer import OpenSearchIndexer
from config.settings import OPENSEARCH_INDEX
from utils.response import success_response

router = APIRouter()

@router.post("/index-folder")
def index_folder(payload: FolderInput):
    client = get_client()
    indexer = OpenSearchIndexer(client, OPENSEARCH_INDEX)
    count = indexer.index_folder(payload.folder)
    return success_response("Folder indexed", {"files_indexed": count})
'''

# =========================
# routes/search_routes.py
# =========================
FILES["routes/search_routes.py"] = r'''
from fastapi import APIRouter
from models.search_models import SearchInput
from opensearch_client.client import get_client
from config.settings import OPENSEARCH_INDEX
from utils.response import success_response

router = APIRouter()

@router.post("/search")
def search(payload: SearchInput):
    client = get_client()

    must = []
    filters = []

    if payload.search_mode == "filename":
        must.append({"match": {"filename": payload.keyword}})
    else:
        must.append({"match": {"content": payload.keyword}})

    if payload.file_types:
        filters.append({"terms": {"filetype": payload.file_types}})

    if payload.date_from or payload.date_to:
        range_q = {}
        if payload.date_from:
            range_q["gte"] = payload.date_from
        if payload.date_to:
            range_q["lte"] = payload.date_to
        filters.append({"range": {"modified": range_q}})

    if payload.size_from or payload.size_to:
        range_q = {}
        if payload.size_from:
            range_q["gte"] = payload.size_from * 1024
        if payload.size_to:
            range_q["lte"] = payload.size_to * 1024
        filters.append({"range": {"size_bytes": range_q}})

    query = {
        "query": {
            "bool": {
                "must": must,
                "filter": filters
            }
        },
        "from": payload.from_,
        "size": payload.size
    }

    res = client.search(index=OPENSEARCH_INDEX, body=query)
    hits = [h["_source"] for h in res["hits"]["hits"]]

    return success_response("Search completed", {
        "count": len(hits),
        "results": hits
    })
'''

# =========================
# WRITE ALL FILES
# =========================
def write_all():
    for path, content in FILES.items():
        full = ROOT / path
        full.parent.mkdir(parents=True, exist_ok=True)
        full.write_text(content.strip() + "\n", encoding="utf-8")
        print("CREATED:", full)

    print("\nProject created at:", ROOT)
    print("Run:")
    print(f"cd {PROJECT_NAME}")
    print("pip install -r requirements.txt")
    print("uvicorn main:app --reload")

if __name__ == "__main__":
    write_all()
